#!/usr/bin/env python3
"""Generate the 3-slide CTB coronagraph-mask deck (ctb_mask_families.ppt).

Slide 1  the basic diffraction model the mask families all hang on
Slide 2  the mask families + code (drivers/helpers/files/params) + compare figure
Slide 3  references

Run:  python3 make_mask_deck.py   (writes ctb_mask_families.ppt in this dir)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(HERE, "ctb_mask_families.ppt")

# ---- palette ---------------------------------------------------------
NAVY   = RGBColor(0x1F, 0x33, 0x55)
BLUE   = RGBColor(0x1A, 0x5B, 0xA8)
INK    = RGBColor(0x20, 0x20, 0x20)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF0, 0xF3, 0xF7)
ACCENT = RGBColor(0xB0, 0x38, 0x0B)
GREEN  = RGBColor(0x1E, 0x6B, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.333)     # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK  = prs.slide_layouts[6]


def add_box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tf


def para(tf, runs, size=14, color=INK, bold=False, bullet=None,
         space_after=4, align=PP_ALIGN.LEFT, first=False, level=0, italic=False):
    """runs = str or list of (text, {opts}) tuples for mixed formatting."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.level = level
    if isinstance(runs, str):
        runs = [(runs, {"italic": italic})]
    if bullet is not None:
        runs = [(bullet + "  ", {})] + list(runs)
    for txt, o in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(o.get("size", size))
        r.font.bold = o.get("bold", bold)
        r.font.color.rgb = o.get("color", color)
        r.font.name = o.get("font", "Calibri")
        if o.get("italic"): r.font.italic = True
    return p


def title_bar(slide, title, sub=None):
    add_box(slide, 0, 0, SW, Inches(0.95), fill=NAVY)
    tf = textbox(slide, Inches(0.4), Inches(0.06), Inches(12.5), Inches(0.85),
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, size=26, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, first=True)
    if sub:
        para(tf, sub, size=13, color=RGBColor(0xC9, 0xD6, 0xE5), bold=False, space_after=0)


def mono(txt):
    return (txt, {"font": "Consolas", "color": BLUE})


# =====================================================================
# SLIDE 1 — the basic diffraction model
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The CTB diffraction model these masks hang on",
          "All-reflective 2-DM coronagraph relay — masks applied to the complex field in MATLAB (no engine change)")

# left column: light path + how a mask is applied
lf = textbox(s, Inches(0.4), Inches(1.15), Inches(6.35), Inches(5.9))
para(lf, "Light path (compact deck ctb_dcr.in, 31 elts)", size=16, bold=True, color=NAVY, first=True, space_after=6)
para(lf, [("source → OAP1 → ", {}), ("DM1", {"bold":True,"color":ACCENT}),
          (" → DM2 → OAP2 → [focus] → OAP3 →", {})], size=13)
para(lf, [("→ ", {}), ("apodizer", {"bold":True,"color":ACCENT}),
          (" (pupil) → OAP4 → ", {}), ("FPM", {"bold":True,"color":ACCENT}),
          (" (focus) → OAP5 →", {})], size=13)
para(lf, [("→ ", {}), ("Lyot", {"bold":True,"color":ACCENT}),
          (" (pupil) → … → ExitPupil → ", {}),
          ("FPA", {"bold":True,"color":ACCENT}), (" (the PSF)", {})], size=13, space_after=12)

para(lf, "The three mask stations", size=16, bold=True, color=NAVY, space_after=6)
para(lf, [("Apodizer", {"bold":True}), (" (pupil) · ", {}),
          ("FPM", {"bold":True}), (" (focus) · ", {}),
          ("Lyot", {"bold":True}), (" (pupil)", {})], size=13, space_after=2)
para(lf, "— passive Reference markers in the deck; masks are added in MATLAB.",
     size=12, color=GREY, space_after=12)

para(lf, "How a mask is applied (stage-1 contract)", size=16, bold=True, color=NAVY, space_after=6)
para(lf, [("propagate to the plane → multiply the complex field by the mask "
           "array → continue.", {})], size=13, bullet="1.", space_after=3)
para(lf, [mono("macos.apodize"), (" (real 0–1 amplitude) · ", {}),
          mono("macos.apodize_complex"), (" (complex).", {})], size=13, bullet="2.", space_after=3)
para(lf, [("subsequent reads use ", {}), mono("'reset_trace',false"),
          (" to keep the mask.", {})], size=13, bullet="3.", space_after=3)
para(lf, "No source tilt, no reference re-alignment, no engine edit.",
     size=12, color=GREEN, bullet="✓", space_after=0)

# right column: the numbers / sampling box
add_box(s, Inches(7.0), Inches(1.2), Inches(5.9), Inches(5.75), fill=LIGHT)
rf = textbox(s, Inches(7.25), Inches(1.35), Inches(5.4), Inches(5.5))
para(rf, "Deterministic sampling (engine geometry, not dx_at at NF planes)",
     size=15, bold=True, color=NAVY, first=True, space_after=8)
para(rf, [("focal pitch   ", {"font":"Consolas"}),
          ("dx_f = λR / (N·dx_sph)", {"font":"Consolas","color":BLUE,"bold":True})],
     size=13, space_after=4)
para(rf, [("FPM-local λ/D   ", {"font":"Consolas"}),
          ("λR / D_beam  (metres)", {"font":"Consolas","color":BLUE,"bold":True})],
     size=13, space_after=4)
para(rf, [("FPA λ/D   ", {"font":"Consolas"}),
          ("λ·R_fpa / D_ep / dx_FPA", {"font":"Consolas","color":BLUE,"bold":True})],
     size=13, space_after=10)
para(rf, "On this bench (N=1024, 500 nm)", size=15, bold=True, color=NAVY, space_after=6)
para(rf, "λ/D at FPA = 4.03 px  (N/pupil; zero-pad N → finer PSF)", size=13, bullet="•", space_after=3)
para(rf, "occulter centred on floor(N/2) = FFT DC pixel (the half-pixel rule)", size=13, bullet="•", space_after=3)
para(rf, "masks re-evaluated PER λ on each λ's grid (chromatic study)", size=13, bullet="•", space_after=3)
para(rf, "2-DM relay → full annular dark zone is the fair scoring region", size=13, bullet="•", space_after=10)
para(rf, "Validation", size=15, bold=True, color=NAVY, space_after=6)
para(rf, "FPM through-focus leg vs MATLAB PROPER: peak-norm corr 1.000000, "
         "dx ratio 1.0000, 0.000 px centroid offset.", size=13, color=GREEN, space_after=0)

# =====================================================================
# SLIDE 2 — mask families + code
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Literature coronagraph mask families",
          "Standard mask/apodizer pairings on the existing apodize machinery — formulae verified verbatim from the papers")

# comparison figure (right)
img = os.path.join(HERE, "ctb_mask_compare.png")
iw = Inches(6.6); ih = Emu(int(iw * 1679 / 3187))
s.shapes.add_picture(img, Inches(6.55), Inches(1.15), width=iw)
cap = textbox(s, Inches(6.55), Inches(1.15) + ih + Inches(0.02), Inches(6.6), Inches(0.3))
para(cap, "ctb_mask_compare.png — contrast vs throughput; APLC deepest, vortex best throughput",
     size=10, color=GREY, italic=True, first=True, align=PP_ALIGN.CENTER)

# table of families (left)
tf = textbox(s, Inches(0.35), Inches(1.1), Inches(6.05), Inches(3.35))
para(tf, "Five families — driver · result (3–15 λ/D, N=1024)", size=14, bold=True, color=NAVY, first=True, space_after=5)
rows = [
 ("Band-limited Lyot", "ctb_bandlimited", "null 4.4e5; (1−ε)R trim verified; 1.1× chromatic", ACCENT),
 ("Vortex + matched Lyot", "ctb_vortex_matched", "3.2× throughput vs unmatched (star→ring)", GREEN),
 ("APLC (prolate)", "ctb_aplc", "174× deeper than BLC at equal 27% thru", BLUE),
 ("Roddier π / dual-zone", "ctb_phase_masks", "first literature complex focal masks", NAVY),
 ("all six, head-to-head", "ctb_mask_compare", "table + scatter + .mat", GREY),
]
for name, drv, res, col in rows:
    para(tf, [(name, {"bold":True, "color":col, "size":13}),
              ("   " + drv, {"font":"Consolas","color":BLUE,"size":11})], space_after=1)
    para(tf, "     " + res, size=11, color=GREY, space_after=4)

# helpers / files / params (left, lower)
add_box(s, Inches(0.35), Inches(4.55), Inches(6.05), Inches(2.55), fill=LIGHT)
bf = textbox(s, Inches(0.55), Inches(4.62), Inches(5.65), Inches(2.45))
para(bf, "Mask builders (new, reusable)", size=12, bold=True, color=NAVY, first=True, space_after=2)
para(bf, [mono("ctb_mask_bandlimited"), (" · ", {}), mono("ctb_apod_prolate"),
          (" · ", {}), mono("ctb_mask_phase")], size=11, space_after=6)
para(bf, "Shared helpers (reused)", size=12, bold=True, color=NAVY, space_after=2)
para(bf, [mono("ctb_mask_disk/softcircle"), (" · ", {}),
          mono("dark_zone_metrics"), (" · ", {}), mono("radial_contrast")], size=11, space_after=6)
para(bf, "Saved (in the example dir)", size=12, bold=True, color=NAVY, space_after=2)
para(bf, "one PNG per driver + ctb_mask_compare.mat (results table) + CTB_PROP_STATUS.md",
     size=11, space_after=6)
para(bf, "User-changeable parameters", size=12, bold=True, color=NAVY, space_after=2)
para(bf, [("model_size, inner/outer_lamD, rx/elt; ", {"font":"Consolas","size":10.5}),
          ("order/form/epsilon", {"font":"Consolas","size":10.5,"color":ACCENT}),
          (" (BLC), ", {"font":"Consolas","size":10.5}),
          ("charges/lyot_fracs", {"font":"Consolas","size":10.5,"color":ACCENT}),
          (" (vortex), ", {"font":"Consolas","size":10.5}),
          ("r_occ_lamD/r_lyot_frac", {"font":"Consolas","size":10.5,"color":ACCENT}),
          (" (APLC), ", {"font":"Consolas","size":10.5}),
          ("rho0_lamD/dz_*", {"font":"Consolas","size":10.5,"color":ACCENT}),
          (" (phase)", {"font":"Consolas","size":10.5})], size=10.5, space_after=0)

# HLC deferral footnote
ft = textbox(s, Inches(6.55), Inches(6.95), Inches(6.6), Inches(0.45))
para(ft, [("HLC deferred to the FALCO integration — its FPM is a "
           "co-optimization product, not a formula.", {"italic":True})],
     size=11, color=ACCENT, first=True)

# =====================================================================
# SLIDE 3 — references
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "References",
          "Formulae extracted verbatim from the ar5iv/arXiv LaTeX of each source")

refs = [
 ("Band-limited Lyot", [
   "Kuchner & Traub 2002, ApJ 570, 900 — 4th-order 1−sinc amplitude mask "
   "(Eq. 7/8); Lyot trim (1−ε).  [astro-ph/0203455]",
   "Kuchner, Crepp & Ge 2005, ApJ 628, 466 — 8th-order masks (Eq. 12/13, m=1,l=3).  [astro-ph/0411077]",
 ]),
 ("APLC (apodized-pupil Lyot)", [
   "Soummer 2005, ApJ 618, L161 — prolate apodizer eigenvalue eqn (Eq. 3).  [astro-ph/0412221]",
   "Soummer, Aime & Falloon 2003, A&A 397, 1161 · Soummer et al. 2011, ApJ 729, 144 (GPI config, 2.8 λ/D).",
 ]),
 ("Phase masks", [
   "Roddier & Roddier 1997, PASP 109, 815 — π spot, radius 0.53 λ/D, no apodizer (flux balance).",
   "N'Diaye et al. 2012, A&A 538, A55 — achromatic dual-zone (pure-phase, wavelength-sliding).  [arXiv:1111.3194]",
   "Soummer, Dohlen & Aime 2003, A&A 403, 369 — dual-zone phase mask.",
 ]),
 ("Vortex", [
   "Mawet et al. 2005, ApJ 633, 1191 · Foo, Palacios & Swartzlander 2005, Opt. Lett. 30, 3308.",
   "Jenkins 2008, MNRAS 384, 515 — ideal even-charge field = 0 inside pupil (Eq. 1).  [arXiv:0709.0153]",
 ]),
]
tf = textbox(s, Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.4))
first = True
for group, items in refs:
    para(tf, group, size=17, bold=True, color=NAVY, first=first, space_after=3)
    first = False
    for it in items:
        para(tf, it, size=13, color=INK, bullet="•", space_after=3)
    para(tf, "", size=6, space_after=2)

nb = textbox(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.5))
para(nb, [("Note: ", {"bold":True, "color":ACCENT}),
          ("a WebFetch summarizer hallucinated equations on every topic; all forms above "
           "were reconciled against the raw paper HTML.", {"italic":True})],
     size=11, color=GREY, first=True)

prs.save(OUT)
print("wrote", OUT)
